"""
Generate a python-pptx Presentation from ACPWB slide data.
Each slide type gets a clean, branded layout matching the web design.
"""
from io import BytesIO

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Brand colors
NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_MID = RGBColor(0x12, 0x20, 0x40)
NAVY_LIGHT = RGBColor(0x1E, 0x35, 0x60)
GOLD = RGBColor(0xC9, 0xA8, 0x4C)
GOLD_LIGHT = RGBColor(0xE0, 0xC0, 0x6E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF6, 0xF9)
MID_GRAY = RGBColor(0xE4, 0xE8, 0xEF)
DARK_TEXT = RGBColor(0x22, 0x22, 0x33)
MID_TEXT = RGBColor(0x55, 0x55, 0x66)

W = Inches(13.333)
H = Inches(7.5)


# ── helpers ─────────────────────────────────────────────────────────────────

def _blank_slide(prs):
    blank_layout = prs.slide_layouts[6]  # index 6 = Blank
    return prs.slides.add_slide(blank_layout)


def _fill_solid(shape, color):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color


def _set_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, left, top, width, height):
    """Add a blank rectangle shape."""
    from pptx.util import Emu
    return slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Emu(left), Emu(top), Emu(width), Emu(height),
    )


def _txbox(slide, left, top, width, height):
    return slide.shapes.add_textbox(
        Emu(left), Emu(top), Emu(width), Emu(height),
    )


def _para(tf, text, size, bold=False, color=WHITE, align=PP_ALIGN.LEFT,
           italic=False, space_before=0, space_after=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return p


def _gold_bar(slide, top_emu, width_emu=None, height_emu=None, left_emu=None):
    """Horizontal gold accent bar."""
    bar = _box(slide,
               left=left_emu or Inches(0),
               top=top_emu,
               width=width_emu or W,
               height=height_emu or Inches(0.045))
    _fill_solid(bar, GOLD)
    bar.line.fill.background()
    return bar


def _header_band(slide, title, height_frac=0.22):
    """Dark navy header band with gold accent + white title text."""
    band_h = int(H * height_frac)
    band = _box(slide, 0, 0, W, band_h)
    _fill_solid(band, NAVY)
    band.line.fill.background()

    # gold accent bar at bottom of header
    _gold_bar(slide, top_emu=band_h - Inches(0.045))

    # title text inside header
    margin = Inches(0.55)
    tb = _txbox(slide, margin, Inches(0.12), W - margin * 2, band_h - Inches(0.2))
    tf = tb.text_frame
    tf.word_wrap = True
    _para(tf, title, size=22, bold=True, color=WHITE)
    return band_h


def _slide_footer(slide, meta, slide_data):
    """Org name + slide number in the bottom strip."""
    footer_h = Inches(0.32)
    footer_top = H - footer_h
    bar = _box(slide, 0, footer_top, W, footer_h)
    _fill_solid(bar, NAVY_MID)
    bar.line.fill.background()

    tb = _txbox(slide, Inches(0.35), footer_top + Inches(0.05), W - Inches(0.7), footer_h)
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = f"{meta.get('org_name', 'ACPWB')}"
    r.font.size = Pt(8)
    r.font.color.rgb = GOLD
    r.font.bold = True

    num_tb = _txbox(slide, W - Inches(1.5), footer_top + Inches(0.05), Inches(1.2), footer_h)
    ntf = num_tb.text_frame
    np_ = ntf.paragraphs[0]
    np_.alignment = PP_ALIGN.RIGHT
    nr = np_.add_run()
    nr.text = f"Slide {slide_data['num']} / {slide_data['total']}"
    nr.font.size = Pt(8)
    nr.font.color.rgb = MID_GRAY


# ── slide type builders ──────────────────────────────────────────────────────

def _slide_title(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY)
    _gold_bar(s, top_emu=Inches(1.5), width_emu=Inches(0.9), left_emu=Inches(0.65), height_emu=Inches(0.055))

    title_tb = _txbox(s, Inches(0.65), Inches(1.65), Inches(9.5), Inches(2.5))
    tf = title_tb.text_frame
    tf.word_wrap = True
    _para(tf, slide.get('heading', meta.get('title', '')), size=34, bold=True, color=WHITE)

    sub_tb = _txbox(s, Inches(0.65), Inches(4.1), Inches(9.5), Inches(1.2))
    stf = sub_tb.text_frame
    stf.word_wrap = True
    _para(stf, slide.get('subheading', meta.get('subtitle', '')), size=16, color=GOLD_LIGHT)

    authors = slide.get('authors', [])
    if authors:
        a = authors[0]
        auth_tb = _txbox(s, Inches(0.65), Inches(5.4), Inches(9.5), Inches(1.2))
        atf = auth_tb.text_frame
        _para(atf, a.get('full_name', ''), size=13, bold=True, color=WHITE)
        _para(atf, a.get('title', ''), size=11, color=GOLD_LIGHT)

    org_tb = _txbox(s, Inches(0.65), Inches(6.6), Inches(9.5), Inches(0.6))
    otf = org_tb.text_frame
    _para(otf, meta.get('org_name', '').upper(), size=11, bold=True, color=GOLD, space_before=0)

    date_str = f"{meta.get('year', '')}–{meta.get('industry', '')}"
    date_tb = _txbox(s, W - Inches(3.5), Inches(6.6), Inches(3.2), Inches(0.6))
    dtf = date_tb.text_frame
    _para(dtf, date_str, size=10, color=MID_GRAY, align=PP_ALIGN.RIGHT)

    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_agenda(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', 'Agenda'))

    items = slide.get('items', [])
    left = Inches(0.75)
    content_top = band_h + Inches(0.35)
    content_w = W - Inches(1.5)
    tb = _txbox(s, left, content_top, content_w, H - band_h - Inches(0.7))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items, 1):
        _para(tf, f"{i}.  {item}", size=14, color=DARK_TEXT, space_before=4)

    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_content(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''))

    bullets = slide.get('bullets', [])
    tb = _txbox(s, Inches(0.75), band_h + Inches(0.35), W - Inches(1.5), H - band_h - Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.size = Pt(13)
        r.font.color.rgb = DARK_TEXT

    _footnote_para(s, slide)
    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_stat(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY)
    heading = slide.get('heading', '')
    _gold_bar(s, top_emu=0, height_emu=Inches(0.05))

    htb = _txbox(s, Inches(0.65), Inches(0.25), W - Inches(1.3), Inches(0.7))
    htf = htb.text_frame
    _para(htf, heading, size=16, bold=True, color=WHITE)

    stats = slide.get('stats', [])
    n = len(stats)
    col_w = (W - Inches(1.0)) / max(n, 1)
    for i, stat in enumerate(stats):
        left = Inches(0.5) + i * col_w
        # stat box
        box = _box(s, left + Inches(0.1), Inches(1.2), col_w - Inches(0.2), Inches(4.8))
        _fill_solid(box, NAVY_MID)
        box.line.fill.background()

        vtb = _txbox(s, left + Inches(0.2), Inches(1.7), col_w - Inches(0.4), Inches(2.0))
        vtf = vtb.text_frame
        _para(vtf, stat.get('value', ''), size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        ltb = _txbox(s, left + Inches(0.2), Inches(3.8), col_w - Inches(0.4), Inches(1.8))
        ltf = ltb.text_frame
        ltf.word_wrap = True
        _para(ltf, stat.get('label', ''), size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_quote(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY_LIGHT)
    _gold_bar(s, top_emu=0, height_emu=Inches(0.06))

    quote = slide.get('quote', '')
    attr = slide.get('attribution', '')

    qtb = _txbox(s, Inches(1.2), Inches(1.2), W - Inches(2.4), Inches(4.2))
    qtf = qtb.text_frame
    qtf.word_wrap = True
    _para(qtf, f"“{quote}”", size=20, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    atb = _txbox(s, Inches(1.2), Inches(5.5), W - Inches(2.4), Inches(0.8))
    atf = atb.text_frame
    _para(atf, f"— {attr}", size=13, color=GOLD, align=PP_ALIGN.CENTER)

    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_chart(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''))

    note_tb = _txbox(s, Inches(0.75), band_h + Inches(0.5), W - Inches(1.5), Inches(3.5))
    ntf = note_tb.text_frame
    ntf.word_wrap = True
    _para(ntf, "[Chart: see interactive version at acpwb.com]", size=13, color=MID_TEXT, italic=True)

    chart_type = slide.get('chart_type', '')
    if chart_type in ('bar_h', 'bar_v'):
        bars = slide.get('chart_bars', [])
        row_top = band_h + Inches(1.2)
        row_h = min(Inches(0.38), (H - band_h - Inches(1.8)) / max(len(bars), 1))
        for bar in bars:
            rtb = _txbox(s, Inches(0.75), row_top, Inches(3.5), row_h)
            rtf = rtb.text_frame
            _para(rtf, bar.get('label', ''), size=10, color=DARK_TEXT)
            vtb = _txbox(s, Inches(4.5), row_top, Inches(2.0), row_h)
            vtf = vtb.text_frame
            _para(vtf, str(bar.get('value', '')), size=10, bold=True, color=DARK_TEXT)
            row_top += row_h + Inches(0.05)
    elif chart_type == 'line':
        pts = slide.get('chart_line_pts', [])
        row_top = band_h + Inches(1.2)
        row_h = Inches(0.35)
        for pt in pts[:8]:
            rtb = _txbox(s, Inches(0.75), row_top, Inches(6.0), row_h)
            rtf = rtb.text_frame
            _para(rtf, f"{pt.get('label', '')}:  {pt.get('value', '')}", size=10, color=DARK_TEXT)
            row_top += row_h

    src = slide.get('chart_source', '')
    if src:
        stb = _txbox(s, Inches(0.75), H - Inches(0.65), W - Inches(1.5), Inches(0.3))
        stf = stb.text_frame
        _para(stf, f"Source: {src}", size=8, color=MID_TEXT, italic=True)

    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_image(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY)

    caption = slide.get('caption', '')
    if caption:
        _gold_bar(s, top_emu=H - Inches(1.2), height_emu=Inches(1.2))
        ctb = _txbox(s, Inches(0.65), H - Inches(1.1), W - Inches(1.3), Inches(1.0))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        _para(ctf, caption, size=12, color=NAVY)

    _slide_footer(s, meta, slide)


def _slide_meme(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY_MID)

    tb = _txbox(s, Inches(1.0), Inches(2.5), W - Inches(2.0), Inches(2.5))
    tf = tb.text_frame
    _para(tf, "[See web version for image]", size=16, color=GOLD, align=PP_ALIGN.CENTER, italic=True)

    _slide_footer(s, meta, slide)


def _slide_summary(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', 'Key Takeaways'))

    bullets = slide.get('bullets', [])
    tb = _txbox(s, Inches(0.75), band_h + Inches(0.4), W - Inches(1.5), H - band_h - Inches(0.8))
    tf = tb.text_frame
    tf.word_wrap = True
    for b in bullets:
        p = tf.add_paragraph()
        p.space_before = Pt(5)
        r = p.add_run()
        r.text = f"•  {b}"
        r.font.size = Pt(13)
        r.font.bold = True
        r.font.color.rgb = DARK_TEXT

    _footnote_para(s, slide)
    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_two_column(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''))

    mid_x = W // 2
    content_top = band_h + Inches(0.2)
    content_h = H - band_h - Inches(0.5)
    col_w = mid_x - Inches(0.75)

    # divider
    div = _box(s, mid_x - Inches(0.01), content_top, Inches(0.02), content_h)
    _fill_solid(div, MID_GRAY)
    div.line.fill.background()

    for side, left_x, label_key, items_key in [
        ('left', Inches(0.55), 'left_label', 'left_items'),
        ('right', mid_x + Inches(0.2), 'right_label', 'right_items'),
    ]:
        label = slide.get(label_key, '')
        items = slide.get(items_key, [])

        ltb = _txbox(s, left_x, content_top + Inches(0.15), col_w, Inches(0.45))
        ltf = ltb.text_frame
        _para(ltf, label, size=12, bold=True, color=GOLD)

        itb = _txbox(s, left_x, content_top + Inches(0.7), col_w, content_h - Inches(0.8))
        itf = itb.text_frame
        itf.word_wrap = True
        for item in items:
            p = itf.add_paragraph()
            p.space_before = Pt(3)
            r = p.add_run()
            r.text = f"•  {item}"
            r.font.size = Pt(11)
            r.font.color.rgb = DARK_TEXT

    _slide_footer(s, meta, slide)
    if slide.get('speaker_note'):
        s.notes_slide.notes_text_frame.text = slide['speaker_note']


def _slide_timeline(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''))

    milestones = slide.get('milestones', [])
    n = max(len(milestones), 1)
    avail_w = W - Inches(1.0)
    col_w = avail_w / n
    box_top = band_h + Inches(0.45)
    box_h = H - band_h - Inches(0.9)

    # timeline connector line
    line_y = box_top + Inches(0.45)
    connector = _box(s, Inches(0.5) + col_w / 2, line_y, avail_w - col_w, Inches(0.04))
    _fill_solid(connector, GOLD)
    connector.line.fill.background()

    for i, ms in enumerate(milestones):
        left = Inches(0.5) + i * col_w
        # dot
        dot = _box(s, left + col_w / 2 - Inches(0.12), line_y - Inches(0.12), Inches(0.24), Inches(0.24))
        _fill_solid(dot, NAVY)
        dot.line.color.rgb = GOLD

        # label (above)
        ltb = _txbox(s, left + Inches(0.1), box_top, col_w - Inches(0.2), Inches(0.38))
        ltf = ltb.text_frame
        ltf.word_wrap = True
        _para(ltf, ms.get('label', ''), size=9, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

        # desc (below)
        dtb = _txbox(s, left + Inches(0.1), line_y + Inches(0.2), col_w - Inches(0.2), box_h - Inches(0.7))
        dtf = dtb.text_frame
        dtf.word_wrap = True
        _para(dtf, ms.get('desc', ''), size=9, color=MID_TEXT, align=PP_ALIGN.CENTER)

    _slide_footer(s, meta, slide)


def _slide_section_divider(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY)
    _gold_bar(s, top_emu=Inches(3.1), width_emu=Inches(1.2), left_emu=Inches(0.65), height_emu=Inches(0.065))

    tb = _txbox(s, Inches(0.65), Inches(3.25), W - Inches(1.3), Inches(2.5))
    tf = tb.text_frame
    tf.word_wrap = True
    _para(tf, slide.get('heading', ''), size=32, bold=True, color=WHITE)

    num_tb = _txbox(s, Inches(0.65), Inches(2.6), Inches(1.5), Inches(0.55))
    ntf = num_tb.text_frame
    _para(ntf, f"Section {slide.get('section_num', '')}",
          size=11, color=GOLD, bold=True)

    _slide_footer(s, meta, slide)


def _slide_process(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''))

    steps = slide.get('steps', [])
    n = max(len(steps), 1)
    avail_w = W - Inches(1.0)
    step_w = avail_w / n
    box_top = band_h + Inches(0.4)
    box_h = H - band_h - Inches(0.85)

    for i, step in enumerate(steps):
        left = Inches(0.5) + i * step_w
        bx = _box(s, left + Inches(0.08), box_top, step_w - Inches(0.16), box_h)
        _fill_solid(bx, NAVY if i % 2 == 0 else NAVY_MID)
        bx.line.fill.background()

        # step number
        numtb = _txbox(s, left + Inches(0.15), box_top + Inches(0.12),
                       step_w - Inches(0.3), Inches(0.5))
        ntf = numtb.text_frame
        _para(ntf, str(i + 1), size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

        # step name
        nametb = _txbox(s, left + Inches(0.1), box_top + Inches(0.65),
                        step_w - Inches(0.2), Inches(0.55))
        nametf = nametb.text_frame
        nametf.word_wrap = True
        _para(nametf, step.get('name', ''), size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        # step desc
        desctb = _txbox(s, left + Inches(0.1), box_top + Inches(1.3),
                        step_w - Inches(0.2), box_h - Inches(1.4))
        desctf = desctb.text_frame
        desctf.word_wrap = True
        _para(desctf, step.get('desc', ''), size=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    _slide_footer(s, meta, slide)


def _slide_case_study(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''))

    org_type = slide.get('org_type', '')
    otb = _txbox(s, Inches(0.65), band_h + Inches(0.2), W - Inches(1.3), Inches(0.35))
    otf = otb.text_frame
    _para(otf, org_type, size=10, italic=True, color=MID_TEXT)

    sections = [
        ('Challenge', slide.get('challenge', ''), NAVY),
        ('Approach', slide.get('approach', ''), NAVY_MID),
        ('Result', slide.get('result', ''), NAVY_LIGHT),
    ]
    section_w = (W - Inches(1.0)) / 3
    for i, (label, content, bg) in enumerate(sections):
        left = Inches(0.5) + i * section_w
        bx = _box(s, left + Inches(0.05), band_h + Inches(0.65),
                  section_w - Inches(0.1), H - band_h - Inches(1.1))
        _fill_solid(bx, bg)
        bx.line.fill.background()

        ltb = _txbox(s, left + Inches(0.15), band_h + Inches(0.75),
                     section_w - Inches(0.3), Inches(0.4))
        ltf = ltb.text_frame
        _para(ltf, label.upper(), size=9, bold=True, color=GOLD)

        ctb = _txbox(s, left + Inches(0.15), band_h + Inches(1.2),
                     section_w - Inches(0.3), H - band_h - Inches(1.7))
        ctf = ctb.text_frame
        ctf.word_wrap = True
        _para(ctf, content, size=10, color=LIGHT_GRAY)

    _slide_footer(s, meta, slide)


def _slide_callout(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY)
    _gold_bar(s, top_emu=0, height_emu=Inches(0.07))

    stat = slide.get('stat', '')
    desc = slide.get('description', '')

    stb = _txbox(s, Inches(0.8), Inches(1.2), W - Inches(1.6), Inches(2.8))
    stf = stb.text_frame
    _para(stf, stat, size=72, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    dtb = _txbox(s, Inches(1.5), Inches(4.2), W - Inches(3.0), Inches(2.2))
    dtf = dtb.text_frame
    dtf.word_wrap = True
    _para(dtf, desc, size=17, color=WHITE, align=PP_ALIGN.CENTER)

    _slide_footer(s, meta, slide)


def _slide_appendix(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, LIGHT_GRAY)
    band_h = _header_band(s, slide.get('heading', ''), height_frac=0.18)

    content = slide.get('content', '')
    tb = _txbox(s, Inches(0.65), band_h + Inches(0.3), W - Inches(1.3),
                H - band_h - Inches(0.75))
    tf = tb.text_frame
    tf.word_wrap = True
    _para(tf, content, size=11, color=DARK_TEXT)

    _slide_footer(s, meta, slide)


def _slide_qanda(prs, meta, slide):
    s = _blank_slide(prs)
    _set_bg(s, NAVY)
    _gold_bar(s, top_emu=0, height_emu=Inches(0.07))

    htb = _txbox(s, Inches(1.0), Inches(1.8), W - Inches(2.0), Inches(2.0))
    htf = htb.text_frame
    _para(htf, slide.get('heading', 'Questions & Discussion'),
          size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    org = slide.get('org_name', meta.get('org_name', ''))
    email = slide.get('contact_email', '')
    itb = _txbox(s, Inches(1.5), Inches(4.0), W - Inches(3.0), Inches(1.5))
    itf = itb.text_frame
    itf.word_wrap = True
    _para(itf, org, size=14, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    if email:
        _para(itf, email, size=12, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    _slide_footer(s, meta, slide)


def _footnote_para(slide_obj, slide_data):
    fn = slide_data.get('footnote')
    if not fn:
        return
    tb = _txbox(slide_obj, Inches(0.55), H - Inches(0.62), W - Inches(1.1), Inches(0.38))
    tf = tb.text_frame
    _para(tf, fn, size=8, color=MID_TEXT, italic=True)


_SLIDE_BUILDERS = {
    'title': _slide_title,
    'agenda': _slide_agenda,
    'content': _slide_content,
    'stat': _slide_stat,
    'quote': _slide_quote,
    'chart': _slide_chart,
    'image': _slide_image,
    'meme': _slide_meme,
    'summary': _slide_summary,
    'two_column': _slide_two_column,
    'timeline': _slide_timeline,
    'section_divider': _slide_section_divider,
    'process': _slide_process,
    'case_study': _slide_case_study,
    'callout': _slide_callout,
    'appendix': _slide_appendix,
    'qanda': _slide_qanda,
}


def generate_pptx_bytes(pres_meta, slides):
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H

    for slide_data in slides:
        slide_type = slide_data.get('type', 'content')
        builder = _SLIDE_BUILDERS.get(slide_type, _slide_content)
        builder(prs, pres_meta, slide_data)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()
